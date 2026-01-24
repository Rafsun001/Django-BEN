import os
import pickle
import pdfplumber
import warnings
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from pinecone import Pinecone, ServerlessSpec
from openai import OpenAI
import re
import tiktoken
import time
from dotenv import load_dotenv

load_dotenv(r"E:\Rafsun\Django Project BEN\projectFolderBEN\.env")
# ------------------ Suppress Warnings ------------------
warnings.filterwarnings("ignore", category=UserWarning, module="pdfminer")

# ------------------ Load API Keys ------------------

openai_api_key = os.getenv("OPENAI_API_KEY")
pinecone_api_key = os.getenv("PINECONE_API_KEY")
pinecone_env = os.getenv("PINECONE_ENV")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME_3")

# ------------------ Initialize Clients ------------------

pc = Pinecone(api_key=pinecone_api_key)

if PINECONE_INDEX_NAME not in pc.list_indexes().names():
    print(f"Creating index: {PINECONE_INDEX_NAME}")
    pc.create_index(
        name=PINECONE_INDEX_NAME,
        dimension=1536,
        metric="cosine",
        spec=ServerlessSpec(cloud="aws", region=pinecone_env)
    )

index = pc.Index(PINECONE_INDEX_NAME)
client = OpenAI(api_key=openai_api_key)

# ------------------ Your PDF List (20 files) ------------------
pdf_paths = [
    r"_Taught AI Academy Policy 12.06.25.docx.pdf",
    r"Behaviour -  Improving_Behaviour_in_Schools_EEF - Evidence_Review.pdf",
    r"Delivering_school_improvement_through_school_to_school_support_GOV UK May2019.pdf",
    r"EEF Teaching and learning strategies  (1).pdf",
    r"guide_to_the_pupil_premium_EEF-_2024.pdf",
    r"Improving LITERACY_GUIDANCE - EEF.pdf",
    r"Maths - Improving-Mathematics-in-KS2.KS3 - EEF-Update.pdf",
    r"OFSTED Improving_20schools.pdf",
    r"Parental_Engagement_-_Evidence_from_Research_and_Practice- EEF .pdf",
    r"Primary Effective-Professional-Development-Guidance-Report EEF.pdf",
    r"Primary implementation a_schools_guide_to_.pdf",
    r"Rosenshine - Principles-of-Insruction.pdf",
    r"Secondary implementation a_schools_guide_to_.pdf",
    r"School inspection handbook - GOV.UK.pdf",
    r"Secondary Effective-Professional-Development-Guidance-Report EEF.pdf",
    r"SEMH - Social_and_Emotional_Learning_Evidence_Review - EEF.pdf",
    r"SEND - special_educational_needs_in_mainstream_schools_guidance_report_EEF.pdf",
    r"Spending guide Pupil premium EEF_evidence_brief.pdf",
    r"Teacher_Feedback_to_Improve_Pupil_Learning - EEF.pdf",
    r"Teacher-quality-recruitment-and-retention-lit-review-EEF.pdf",
    r"Sample inputs with responses .pdf",
]

# Safety check: warn if any file is missing
print(f"Found {len(pdf_paths)} files in pdf_paths")
for i, path in enumerate(pdf_paths, 1):
    if not os.path.exists(path):
        print(f"   WARNING: File {i} does NOT exist: {path}")

PROCESSED_DATA_FILE = "hs_new_data.pkl"
MAX_HISTORY_LENGTH = 10
MAX_TOTAL_TOKENS = 12000
ENCODING = tiktoken.get_encoding("cl100k_base")

# ------------------ Token Tools ------------------
def num_tokens(text: str) -> int:
    return len(ENCODING.encode(text, disallowed_special=()))

def truncate_by_tokens(text: str, max_tokens: int) -> str:
    tokens = ENCODING.encode(text, disallowed_special=())
    if len(tokens) <= max_tokens:
        return text
    return ENCODING.decode(tokens[:max_tokens])

# ------------------ File Extraction ------------------
def extract_text_from_pdf(path):
    text = ""
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"
    except Exception as e:
        print(f"Error reading PDF {path}: {e}")
    return text

def extract_text_from_docx(path):
    text = ""
    try:
        doc = DocxDocument(path)
        for p in doc.paragraphs:
            text += p.text + "\n"
    except Exception as e:
        print(f"Error reading DOCX {path}: {e}")
    return text

def extract_text_from_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX {path}: {e}")
    return text

def extract_text_from_file(path):
    if path.endswith(".pdf"):
        return extract_text_from_pdf(path)
    elif path.endswith(".docx"):
        return extract_text_from_docx(path)
    elif path.endswith(".pptx"):
        return extract_text_from_pptx(path)
    return ""

def split_chunks(text, size=1000, overlap=200):
    splitter = RecursiveCharacterTextSplitter(chunk_size=size, chunk_overlap=overlap)
    return splitter.split_text(text)

# ------------------ SMART UPLOAD: Only Missing Files ------------------
def smart_upload_knowledge_base():
    if not os.path.exists(PROCESSED_DATA_FILE):
        print("First run → processing all documents...")
        data = {}
        for p in pdf_paths:
            if os.path.exists(p):
                name = os.path.basename(p)
                print(f"   → {name}")
                text = extract_text_from_file(p)
                chunks = split_chunks(text)
                data[name] = chunks
        with open(PROCESSED_DATA_FILE, "wb") as f:
            pickle.dump(data, f)
    else:
        print("Loading cached document chunks...")
        with open(PROCESSED_DATA_FILE, "rb") as f:
            data = pickle.load(f)

    embedder = OpenAIEmbeddings(api_key=openai_api_key)
    missing = []

    print("\nChecking Pinecone for existing sources...")
    for filename in data.keys():
        try:
            dummy_vec = [0.0] * 1535 + [1e-6]
            res = index.query(
                vector=dummy_vec,
                top_k=1,
                filter={"source": {"$eq": filename}},
                include_metadata=False
            )
            if len(res.matches) == 0:
                missing.append(filename)
            else:
                print(f"   Already in Pinecone: {filename}")
        except Exception as e:
            print(f"   Error checking {filename}: {e} → will upload")
            missing.append(filename)

    if not missing:
        print("All documents already in Pinecone! Skipping upload.\n")
        return

    print(f"\n{len(missing)} new/missing document(s) → uploading...")
    for f in missing:
        print(f"   → {f}")

    vectors = []
    total = sum(len(data[f]) for f in missing)
    done = 0

    for filename in missing:
        for i, chunk in enumerate(data[filename]):
            if len(chunk.strip()) < 30:
                continue
            for attempt in range(3):
                try:
                    vec = embedder.embed_query(chunk)
                    vectors.append((
                        f"{filename}_chunk_{i}",
                        vec,
                        {"text": chunk, "source": filename}
                    ))
                    done += 1
                    print(f"\r   Progress: {done}/{total} chunks", end="")
                    break
                except Exception as e:
                    if attempt == 2:
                        print(f"\n   Failed chunk {i}: {e}")
                    else:
                        time.sleep(2 ** attempt)
            if len(vectors) >= 100:
                index.upsert(vectors)
                vectors = []
                time.sleep(0.2)

    if vectors:
        index.upsert(vectors)

    print(f"\nUpload complete! {done} new chunks added.\n")

# ------------------ Clean Output ------------------
def clean(text):
    return re.sub(r'(\*\*|##+|```)', '', text).strip()

# ------------------ MAIN STRATEGIC LEADERSHIP GENERATOR ------------------
def generate_hs_response(
    user_input: str,
    uploaded_files=None,
    history=None
):
    if history is None:
        history = []
    if uploaded_files is None:
        uploaded_files = []

    # 1. Uploaded files
    uploaded_ctx = ""
    for path in uploaded_files:
        if os.path.exists(path):
            text = extract_text_from_file(path)
            if text.strip():
                uploaded_ctx += f"\n\n--- UPLOADED: {os.path.basename(path)} ---\n{text}"

    # 2. Retrieve from Pinecone
    embedder = OpenAIEmbeddings(api_key=openai_api_key)
    query_vec = embedder.embed_query(user_input)
    results = index.query(vector=query_vec, top_k=6, include_metadata=True)

    pinecone_ctx = "\n\n".join([
        f"[Source: {m.metadata.get('source', 'unknown')}]\n{m.metadata.get('text', '')}"
        for m in results.matches
    ])

    full_context = uploaded_ctx + "\n\n" + pinecone_ctx

    system_prompt = """
🏛️ Taught AI Academy – Strategic Leadership GPT

🔐 Acceptable Use Policy

Confirm all input is anonymised and professional.

➡️ If confirmed, continue.

—

🎯 R – Role

You support senior leaders in generating:
- Briefings
- Position statements
- Policy overviews
- CPD drafts
- Evidence-based recommendations

—

📝 I – Instruction

“Before generating a response, refer to the attached ‘Sample Inputs and Responses’ document. Match the tone, style, and structure of the examples provided there as closely as possible.”

Generate:
- 3 evidence-aligned options or strategies
- One polished draft (briefing, summary, or policy outline)
- Ask whether the user would like any of the repsonse expanded on or added to tables 

—

📋 Format

Use:
- Section headers (e.g., Rationale, Impact, Strategy)
- Tables for implementation plans
- Footnotes with policy links or guidance sources

—

📌 References

Align with:
- Ofsted EIF
- DfE statutory guidance
- EEF recommendations
- Taught AI Academy policies

—

🛑 Safeguards

No personal data. No pupil-facing content. Output must be suitable for SLT/Governance use.
"""

    user_message = f"""
USER REQUEST:
{user_input}

REFERENCE DOCUMENTS:
{full_context if full_context.strip() else "No relevant documents found."}
"""

    messages = [
        {"role": "system", "content": system_prompt},
        *history[-MAX_HISTORY_LENGTH:]
    ]

    used = sum(num_tokens(m["content"]) for m in messages)
    available = MAX_TOTAL_TOKENS - used - 2500
    safe_msg = truncate_by_tokens(user_message, available)

    messages.append({"role": "user", "content": safe_msg})

    print(f"Final prompt: ~{sum(num_tokens(m['content']) for m in messages)} tokens")

    response = client.chat.completions.create(
        model="gpt-4-turbo",
        messages=messages,
        temperature=0.7,
        max_tokens=3000
    )

    reply = clean(response.choices[0].message.content)

    history.append({"role": "user", "content": user_input})
    history.append({"role": "assistant", "content": reply})

    return reply

# ------------------ RUN ------------------
if __name__ == "__main__":
    print("Taught AI Academy – Heads & SLT Strategic GPT (Smart Upload Edition)")
    smart_upload_knowledge_base()

    uploaded_files = []  # Add paths here if needed

    query = "Draft a new mobile phone policy for the school following repeated incidents of bullying linked to phone use. Please ensure it aligns with DfE guidance, promotes pupil safety, and includes clear expectations, sanctions, and communication plans for staff, students, and parents."

    result = generate_hs_response(query, uploaded_files)
    print("\n" + "="*80)
    print("STRATEGIC RESPONSE")
    print("="*80)
    print(result)