import os
import pickle
import pdfplumber
import warnings
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from pinecone import Pinecone, ServerlessSpec
from openai import OpenAI
import re
import tiktoken
import time
from dotenv import load_dotenv

load_dotenv(r"E:\Rafsun\Django Project BEN\projectFolderBEN\.env")
# ------------------ Suppress Warnings ------------------
warnings.filterwarnings("ignore", category=UserWarning, module="pdfminer")

# ------------------ Load API Keys ------------------

openai_api_key = os.getenv("OPENAI_API_KEY")
pinecone_api_key = os.getenv("PINECONE_API_KEY")
pinecone_env = os.getenv("PINECONE_ENV")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME_2")

# ------------------ Initialize Clients ------------------

pc = Pinecone(api_key=pinecone_api_key)

if PINECONE_INDEX_NAME not in pc.list_indexes().names():
    print(f"Creating new Pinecone index: {PINECONE_INDEX_NAME}")
    pc.create_index(
        name=PINECONE_INDEX_NAME,
        dimension=1536,
        metric="cosine",
        spec=ServerlessSpec(cloud="aws", region=pinecone_env)
    )

index = pc.Index(PINECONE_INDEX_NAME)
client = OpenAI(api_key=openai_api_key)

# ------------------ Your PDF List ------------------
pdf_paths = [
    r"ADHD - Teaching-and-Managing-Students-with-ADHD - SHIRE.pdf",
    r"ADHD Evidence-based strategies for ADHD.docx.pdf",
    r"ASD school stress and anxiety - Autism Education Trust.pdf",
    r"Autism and Challenging Behaviours_ Strategies and support - Autism Speaks.pdf",
    r"EAL - Classroom Strategies Working with EAL Learners in Secondary Settings - Bell Foundation.pdf",
    r"Improving LITERACY_GUIDANCE - EEF.pdf",
    r"Literacy improvement in KS2 report Second edition EEF.pdf",
    r"Literacy improvements in KS1 Guidance Report 2020 EEF.pdf",
    r"Maths - Improving-Mathematics-in-KS2.KS3 - EEF-Update.pdf",
    r"Maths Early Years to KS1 Guidance Report EEF.pdf",
    r"Maths KS2 KS3 improvements 2022 Guidence report Update EEF.pdf",
    r"Metacognition_and_self-regulated_learning - EEF - .pdf",
    r"Rosenshine - Principles-of-Insruction.pdf",
    r"SEMH - Social_and_Emotional_Learning_Evidence_Review - EEF.pdf",
    r"SEND - special_educational_needs_in_mainstream_schools_guidance_report_EEF.pdf",
    r"Science improvement in primary guidance-report-ks1-ks2 EEF.pdf",
    r"Science improvement strategies EEF .pdf",
    r"TRAUMA-INFORMED INSTRUCTION - Hanover institute.pdf",
    r"Teacher_Feedback_to_Improve_Pupil_Learning - EEF.pdf",
    r"_Taught AI Academy Policy 12.06.25.pdf",
    r"Sample inputs and Responses .pdf"
]

PROCESSED_DATA_FILE = "sow_new_data.pkl"
MAX_HISTORY_LENGTH = 10
MAX_TOTAL_TOKENS = 12000
ENCODING = tiktoken.get_encoding("cl100k_base")

# ------------------ Token Utilities ------------------
def num_tokens_from_string(string: str) -> int:
    return len(ENCODING.encode(string, disallowed_special=()))

def truncate_text_by_tokens(text: str, max_tokens: int) -> str:
    tokens = ENCODING.encode(text, disallowed_special=())
    if len(tokens) <= max_tokens:
        return text
    return ENCODING.decode(tokens[:max_tokens])

# ------------------ File Extraction ------------------
def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"
    except Exception as e:
        print(f"Error reading {pdf_path}: {e}")
    return text

def extract_text_from_docx(docx_path):
    text = ""
    try:
        doc = DocxDocument(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error reading DOCX {docx_path}: {e}")
    return text

def extract_text_from_pptx(pptx_path):
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX {pptx_path}: {e}")
    return text

def extract_text_from_file(file_path):
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    return ""

def split_text_into_chunks(text, chunk_size=1000, chunk_overlap=200):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_text(text)

# ------------------ SMART PROCESS & UPLOAD (Only Missing Files) ------------------
def process_and_upload_smart():
    # 1. Load / create processed chunks
    if not os.path.exists(PROCESSED_DATA_FILE):
        print("First run → extracting and caching all PDFs...")
        data = {}
        for path in pdf_paths:
            if os.path.exists(path):
                name = os.path.basename(path)
                print(f"   → Extracting: {name}")
                text = extract_text_from_file(path)
                chunks = split_text_into_chunks(text)
                data[name] = chunks
        with open(PROCESSED_DATA_FILE, "wb") as f:
            pickle.dump(data, f)
    else:
        print("Loading cached PDF chunks from disk...")
        with open(PROCESSED_DATA_FILE, "rb") as f:
            data = pickle.load(f)

    # 2. Find which SOURCES are really missing in Pinecone
    embeddings_model = OpenAIEmbeddings(api_key=openai_api_key)
    missing_files = []

    print("\nChecking Pinecone for existing sources...")
    for filename in data.keys():
        try:
            # Use a tiny random vector so Pinecone always returns something if the filter matches
            dummy_vector = [0.0] * 1535 + [0.0001]   # not all zeros
            results = index.query(
                vector=dummy_vector,
                top_k=1,
                filter={"source": {"$eq": filename}},
                include_metadata=False
            )
            if len(results.matches) == 0:
                missing_files.append(filename)
            else:
                print(f"   Already present: {filename}")
        except Exception as e:
            print(f"   Error checking {filename}: {e} → will re-upload")
            missing_files.append(filename)

    if not missing_files:
        print("All files are already in Pinecone! Nothing to do.")
        return

    print(f"\n{len(missing_files)} file(s) need uploading:")
    for f in missing_files:
        print(f"   → {f}")

    # 3. Upload only the missing ones (with proper batching + retries)
    vectors = []
    total_chunks = sum(len(data[f]) for f in missing_files)
    uploaded = 0

    for filename in missing_files:
        chunks = data[filename]
        for i, chunk in enumerate(chunks):
            if len(chunk.strip()) < 30:
                continue

            for attempt in range(3):  # retry up to 3 times
                try:
                    emb = embeddings_model.embed_query(chunk)
                    vectors.append((
                        f"{filename}_chunk_{i}",
                        emb,
                        {"text": chunk, "source": filename}
                    ))
                    uploaded += 1
                    print(f"\r   Embedded: {uploaded}/{total_chunks} chunks", end="")
                    break  # success → exit retry loop
                except Exception as e:
                    print(f"\n   Retry {attempt+1}/3 for chunk {i} of {filename}: {e}")
                    time.sleep(2 ** attempt)  # exponential back-off
            else:
                print(f"\n   Failed to embed chunk {i} of {filename}")

            # Batch upsert every 100 vectors
            if len(vectors) >= 100:
                try:
                    index.upsert(vectors=vectors)
                    vectors = []
                    time.sleep(0.2)
                except Exception as e:
                    print(f"\n   Upsert failed: {e}")

    # Final batch
    if vectors:
        try:
            index.upsert(vectors=vectors)
        except Exception as e:
            print(f"\n   Final upsert failed: {e}")

    print(f"\nSUCCESS! Uploaded {uploaded} new chunks.")
    print("Your knowledge base is now 100% up to date.\n")

# ------------------ Clean Output ------------------
def remove_markdown_junk(text):
    text = re.sub(r'(\*\*|##+|```)', '', text)
    return text.strip()

# ------------------ MAIN CHAT FUNCTION (unchanged logic) ------------------
def get_support_response_sow_academy_new(
    user_input: str,
    uploaded_files=None,
    history=None
):
    if history is None:
        history = []
    if uploaded_files is None:
        uploaded_files = []

    # Uploaded user files
    uploaded_context = ""
    for file_path in uploaded_files:
        if os.path.exists(file_path):
            text = extract_text_from_file(file_path)
            if text.strip():
                uploaded_context += f"\n\n--- UPLOADED: {os.path.basename(file_path)} ---\n{text}"

    # Retrieve from Pinecone
    embeddings = OpenAIEmbeddings(api_key=openai_api_key)
    query_emb = embeddings.embed_query(user_input)
    results = index.query(vector=query_emb, top_k=6, include_metadata=True)

    retrieved_context = "\n\n".join([
        f"[Source: {m.metadata.get('source', 'unknown')}]\n{m.metadata.get('text', '')}"
        for m in results.matches
    ])

    full_context = uploaded_context + "\n\n" + retrieved_context

    user_message = f"""
User request:
"{user_input}"

Reference materials:
{full_context if full_context.strip() else "No relevant documents found."}
"""

    system_message = """
📅 Taught AI Academy – Scheme of Work Generator GPT

🔐 Acceptable Use Check

Ask: “Does your prompt avoid names or personal data?”  
➡️ Proceed only with confirmation.

—

🎯 R – Role

You generate 6–12 week trauma-informed schemes of work for KS3/4 PRU learners.

—

📝 I – Instruction
Before generating a response, refer to the attached ‘Sample Inputs and Responses’ document. Match the tone, style, and structure of the examples provided there as closely as possible.

Ask:
- Subject?
- Theme/topic?
- Accreditation (GCSE/FS/Unit)?
- Barriers to learning?
- Prior knowledge?

Then generate:
- Weekly grid (Week, Topic, LO, Activity, Assessment, Resources)
- Differentiated outcomes
- Key vocabulary
- Behaviour/inclusion support

—

📋 Output Format

Use table layout with bold column headers.  
Avoid markdown formatting. Use trauma-informed vocabulary and adjust for low literacy.

—

💬 Follow-Up Prompt

“Would you like this exported to Word or turned into CPD?”

—

📌 Safeguards

Do not generate for named individuals. No clinical summaries or EHCP-level data.

"""

    messages = [
        {"role": "system", "content": system_message},
        *history[-MAX_HISTORY_LENGTH:]
    ]

    tokens_used = sum(num_tokens_from_string(m["content"]) for m in messages)
    available = MAX_TOTAL_TOKENS - tokens_used - 2500
    safe_user_message = truncate_text_by_tokens(user_message, available)

    messages.append({"role": "user", "content": safe_user_message})

    print(f"Final prompt: ~{sum(num_tokens_from_string(m['content']) for m in messages)} tokens")

    response = client.chat.completions.create(
        model="gpt-4-turbo",
        messages=messages,
        temperature=0.8,
        max_tokens=2500
    )

    reply = response.choices[0].message.content.strip()
    reply = remove_markdown_junk(reply)

    history.append({"role": "user", "content": user_input})
    history.append({"role": "assistant", "content": reply})

    return reply

# ------------------ RUN ONCE ------------------
if __name__ == "__main__":
    print("Taught AI Academy – Scheme of Work Generator (Smart Upload Edition)")
    process_and_upload_smart()  # ← This is the smart version!

    uploaded_files = []  # Add paths here if needed
    user_query = "Create a 6-week Gaming & Mental Health scheme for Year 10 boys in PRU who refuse to write long answers. Use project-based tasks, TikTok-style outcomes, and build in movement breaks."

    response = get_support_response_sow_academy(user_query, uploaded_files)
    print("\n" + "="*70)
    print("SCHEME OF WORK")
    print("="*70)
    print(response)